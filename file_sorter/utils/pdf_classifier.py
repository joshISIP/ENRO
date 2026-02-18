import os
import PyPDF2
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation


# ---------------- TEXT EXTRACTORS ---------------- #

def extract_text_from_pdf(path):
    text = ""
    try:
        with open(path, "rb") as file:
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                text += page.extract_text() or ""
    except:
        pass
    return text.lower()


def extract_text_from_docx(path):
    text = ""
    try:
        doc = Document(path)
        for para in doc.paragraphs:
            text += para.text + " "
    except:
        pass
    return text.lower()


def extract_text_from_xlsx(path):
    text = ""
    try:
        wb = load_workbook(path, data_only=True)
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                for cell in row:
                    if cell:
                        text += str(cell) + " "
    except:
        pass
    return text.lower()


def extract_text_from_pptx(path):
    text = ""
    try:
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + " "
    except:
        pass
    return text.lower()


# ---------------- MAIN CLASSIFIER ---------------- #

def extract_text(file_path):
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext == ".docx":
        return extract_text_from_docx(file_path)
    elif ext == ".xlsx":
        return extract_text_from_xlsx(file_path)
    elif ext == ".pptx":
        return extract_text_from_pptx(file_path)
    else:
        return ""


def classify_letter(file_path):
    text = extract_text(file_path)

    if any(word in text for word in ["invoice", "amount due", "billing", "payment due"]):
        return "Invoice"

    elif any(word in text for word in ["curriculum vitae", "resume", "education", "skills"]):
        return "Resume"

    elif any(word in text for word in ["contract", "agreement", "terms and conditions"]):
        return "Contract"

    elif any(word in text for word in ["request letter", "i am writing to request", "approval"]):
        return "Request Letter"

    elif any(word in text for word in ["report", "summary report", "findings"]):
        return "Report"

    else:
        return "Uncategorized"
